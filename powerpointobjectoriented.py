from pptx import Presentation
from pptx.util import Inches
from Lyric_getter import lyric_getter
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
class Slides:
	def __init__(self,master=None):
		self.pres = Presentation(master)
		
		
		

	def add_text(self,url):
		self.url = url
		self.lyrics = lyric_getter(self.url)
		for paragraph in self.lyrics:
			#self.slide_layout = self.pres.slide_layouts[0]
			#self.slide = self.pres.slides.add_slide(self.slide_layout)
			#self.title = self.slide.shapes.title
			#self.title.text=paragraph 
			#font=self.title.font
			slide=self.pres.slides.add_slide(self.pres.slide_layouts[5])
			shape=slide.shapes
			left = Inches(0)
			top = Inches(2)
			width = Inches(10)
			height = Inches(4)
			shape = shape.add_shape(
    		MSO_SHAPE.RECTANGLE, left, top, width, height
			)
			line = shape.line
			line.fill.background()
			shape.fill.background()
			shape.shadow.inherit = False
			text_frame=shape.text_frame

			
			text_frame.clear()  # not necessary for newly-created shape
			text_frame.margin_bottom = Inches(1)
			text_frame.margin_left = 0
			text_frame.vertical_anchor = MSO_ANCHOR.TOP

			p = text_frame.paragraphs[0]
			run = p.add_run()
			run.text = paragraph
			font=run.font
			font.size=Pt(34)
			font.name = 'Roboto'
			font.color.rgb = RGBColor(100, 100, 100)


	

	def savepresentation(self,name):
		self.pres.save(name)

# a = Slides()
# a.add_text('https://www.azlyrics.com/lyrics/littlemix/countingstarsholygrailsmellsliketeenspirit.html')
# a.savepresentation('objectpowers6.pptx')



